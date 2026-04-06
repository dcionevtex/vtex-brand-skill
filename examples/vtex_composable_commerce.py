"""
VTEX Composable Commerce — Presentation Generator
Audience: CTOs / Enterprise decision-makers
Brand: VTEX official guidelines (brand.vtex.com)

Usage:
    pip install python-pptx
    python3 vtex_composable_commerce.py
    python3 vtex_composable_commerce.py --logo-rebel /path/to/vtex-logo-rebel.png --logo-white /path/to/vtex-logo-white.png
    python3 vtex_composable_commerce.py --output my_deck.pptx

Logo download:
    https://brand.vtex.com/wp-content/uploads/2024/08/Logo-Tagline-Composable-and-Complete.zip
"""

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand Colors
REBEL_PINK    = RGBColor(0xF7, 0x19, 0x63)
SERIOUS_BLACK = RGBColor(0x14, 0x20, 0x32)
PLAIN_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SOFT_BLUE     = RGBColor(0xF5, 0xF9, 0xFF)
SOFT_PINK     = RGBColor(0xFF, 0xF3, 0xF6)
WINTER_GRAY   = RGBColor(0xE7, 0xE9, 0xEE)
COOL_GRAY     = RGBColor(0xA1, 0xA8, 0xB7)
SERIOUS_GRAY  = RGBColor(0x5B, 0x6E, 0x84)

FONT     = 'VTEX Trust'
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)
ML       = Inches(0.8)
CW       = SLIDE_W - Inches(1.6)


def new_deck():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def text(slide, content, left, top, width, height,
         size=20, color=SERIOUS_BLACK, bold=False, italic=False, align='left'):
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(content.split('\n')):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align_map.get(align, PP_ALIGN.LEFT)
        run = para.add_run()
        run.text = line
        run.font.name   = FONT
        run.font.size   = Pt(size)
        run.font.color.rgb = color
        run.font.bold   = bold
        run.font.italic = italic
    return box


def logo(slide, logo_path=None, left=Inches(0.35), top=Inches(0.18), height=Inches(0.48)):
    if logo_path:
        slide.shapes.add_picture(logo_path, left, top, height=height)
    else:
        box = slide.shapes.add_textbox(left, top, Inches(1.3), height)
        run = box.text_frame.paragraphs[0].add_run()
        run.text = 'VTEX'
        run.font.name  = FONT
        run.font.size  = Pt(22)
        run.font.bold  = True
        run.font.color.rgb = REBEL_PINK


def accent_bar(slide, position='top', orientation='h', thickness=Inches(0.045)):
    if orientation == 'h':
        w, h = SLIDE_W, thickness
        l = Inches(0)
        t = Inches(0) if position == 'top' else SLIDE_H - thickness
    else:
        w, h = thickness, SLIDE_H
        l, t = Inches(0), Inches(0)
    return rect(slide, l, t, w, h, REBEL_PINK)


def slide_cover(prs, title, subtitle=None, logo_path=None):
    s = blank_slide(prs)
    bg(s, SERIOUS_BLACK)
    accent_bar(s, position='bottom')
    logo(s, logo_path)
    text(s, title, left=ML, top=Inches(2.0), width=CW, height=Inches(2.4), size=46, color=PLAIN_WHITE)
    if subtitle:
        text(s, subtitle, left=ML, top=Inches(4.6), width=CW, height=Inches(0.8), size=19, color=COOL_GRAY)
    return s


def slide_section(prs, title, eyebrow=None, logo_path=None):
    s = blank_slide(prs)
    bg(s, REBEL_PINK)
    logo(s, logo_path)
    if eyebrow:
        text(s, eyebrow.upper(), left=ML, top=Inches(2.4), width=CW, height=Inches(0.5), size=13, color=PLAIN_WHITE)
    text(s, title, left=ML, top=Inches(2.9), width=CW, height=Inches(2.0), size=42, color=PLAIN_WHITE)
    return s


def slide_content(prs, title, bullets, logo_path=None, note=None):
    s = blank_slide(prs)
    bg(s, SOFT_BLUE)
    accent_bar(s, position='top')
    logo(s, logo_path)
    text(s, title, left=ML, top=Inches(0.9), width=CW, height=Inches(0.85), size=30, color=SERIOUS_BLACK)
    body = '\n'.join(f'\u2013  {b}' for b in bullets)
    text(s, body, left=ML, top=Inches(1.95), width=CW, height=Inches(4.6), size=17, color=SERIOUS_GRAY)
    if note:
        text(s, note, left=ML, top=Inches(6.7), width=CW, height=Inches(0.4), size=12, color=COOL_GRAY, italic=True)
    return s


def slide_two_col(prs, title, left_head, left_items, right_head, right_items, logo_path=None):
    s = blank_slide(prs)
    bg(s, PLAIN_WHITE)
    accent_bar(s, position='top')
    logo(s, logo_path)
    col_w, col_gap = Inches(5.5), Inches(0.5)
    col_l, col_r = ML, ML + col_w + col_gap
    text(s, title, left=ML, top=Inches(0.9), width=CW, height=Inches(0.85), size=30, color=SERIOUS_BLACK)
    rect(s, col_l, Inches(2.0), col_w, Inches(0.55), REBEL_PINK)
    text(s, left_head, left=col_l+Inches(0.15), top=Inches(2.05), width=col_w-Inches(0.3), height=Inches(0.45), size=14, color=PLAIN_WHITE, bold=True)
    text(s, '\n'.join(f'\u2013  {i}' for i in left_items), left=col_l, top=Inches(2.65), width=col_w, height=Inches(4.0), size=16, color=SERIOUS_GRAY)
    rect(s, col_r, Inches(2.0), col_w, Inches(0.55), SERIOUS_BLACK)
    text(s, right_head, left=col_r+Inches(0.15), top=Inches(2.05), width=col_w-Inches(0.3), height=Inches(0.45), size=14, color=PLAIN_WHITE, bold=True)
    text(s, '\n'.join(f'\u2013  {i}' for i in right_items), left=col_r, top=Inches(2.65), width=col_w, height=Inches(4.0), size=16, color=SERIOUS_GRAY)
    return s


def slide_stat(prs, stats, title=None, logo_path=None):
    s = blank_slide(prs)
    bg(s, SERIOUS_BLACK)
    accent_bar(s, position='top')
    logo(s, logo_path)
    if title:
        text(s, title, left=ML, top=Inches(0.9), width=CW, height=Inches(0.7), size=22, color=COOL_GRAY)
    col_w = CW / max(len(stats), 1)
    for i, (number, label) in enumerate(stats):
        x = ML + col_w * i
        text(s, number, left=x, top=Inches(2.2), width=col_w, height=Inches(1.8), size=72, color=REBEL_PINK, align='center')
        text(s, label,  left=x, top=Inches(4.1), width=col_w, height=Inches(0.8), size=16, color=COOL_GRAY,  align='center')
        if i > 0:
            rect(s, x - Inches(0.02), Inches(2.2), Inches(0.02), Inches(2.6), SERIOUS_GRAY)
    return s


def slide_quote(prs, quote, attribution=None, logo_path=None):
    s = blank_slide(prs)
    bg(s, SOFT_PINK)
    logo(s, logo_path)
    rect(s, Inches(0), Inches(0), Inches(0.12), SLIDE_H, REBEL_PINK)
    text(s, f'"{quote}"', left=Inches(1.0), top=Inches(1.8), width=Inches(11.0), height=Inches(3.5), size=26, color=SERIOUS_BLACK, italic=True)
    if attribution:
        text(s, f'\u2014 {attribution}', left=Inches(1.0), top=Inches(5.5), width=Inches(8.0), height=Inches(0.6), size=15, color=SERIOUS_GRAY)
    return s


def slide_closing(prs, headline, cta=None, logo_path=None):
    s = blank_slide(prs)
    bg(s, SERIOUS_BLACK)
    accent_bar(s, orientation='v', thickness=Inches(0.14))
    logo(s, logo_path, left=Inches(0.35), top=SLIDE_H - Inches(1.05), height=Inches(0.48))
    text(s, headline, left=Inches(1.1), top=Inches(2.2), width=Inches(10.5), height=Inches(2.0), size=42, color=PLAIN_WHITE)
    if cta:
        text(s, cta, left=Inches(1.1), top=Inches(4.4), width=Inches(7.0), height=Inches(0.7), size=18, color=REBEL_PINK)
    return s


def build_deck(output_path, logo_rebel=None, logo_white=None):
    prs = new_deck()

    slide_cover(prs, "Composable Commerce.\nBuilt for how you actually operate.",
                subtitle="Enterprise Architecture — VTEX", logo_path=logo_white)
    slide_section(prs, "Monolithic platforms are slowing you down.",
                  eyebrow="Section 01", logo_path=logo_white)
    slide_content(prs, "Where the friction shows up", bullets=[
        "New features take quarters, not weeks — the platform owns your roadmap",
        "Frontend, backend, and infrastructure locked together — one change breaks three things",
        "Scaling for peak traffic means scaling everything, not just what matters",
        "Integrating new channels requires deep customization with every release",
        "Engineering spends 60%+ of time maintaining infrastructure, not building",
    ], logo_path=logo_rebel)
    slide_stat(prs, title="The cost of staying still", stats=[
        ("18mo", "Average release cycle\non monolithic platforms"),
        ("3\u00d7",  "Longer time-to-market\nvs composable peers"),
        ("40%",  "Revenue impact from\npoor scalability on peak days"),
    ], logo_path=logo_white)
    slide_section(prs, "What composable actually means.", eyebrow="Section 02", logo_path=logo_white)
    slide_content(prs, "Composable is not just headless", bullets=[
        "Headless = decoupled frontend. Composable = every capability independently replaceable",
        "Built on MACH principles: Microservices, API-first, Cloud-native, Headless",
        "Each commerce function — checkout, catalog, promotions, OMS — is a replaceable module",
        "You keep what works and swap what doesn't, without a full re-platform",
        "The result: your platform evolves with your business, not against it",
    ], note="MACH = Microservices \u00b7 API-first \u00b7 Cloud-native \u00b7 Headless", logo_path=logo_rebel)
    slide_two_col(prs, "Monolith vs Composable: what changes",
                  left_head="Monolithic", left_items=[
                      "One release cycle for everything",
                      "Coupled frontend and backend",
                      "Customization via platform plugins",
                      "Scale the whole stack or none of it",
                      "Vendor owns the roadmap",
                  ],
                  right_head="Composable", right_items=[
                      "Deploy each capability independently",
                      "Frontend teams move at their own pace",
                      "Extend via APIs — any language, any tool",
                      "Scale only what needs to scale",
                      "You own the decisions",
                  ], logo_path=logo_rebel)
    slide_section(prs, "How VTEX enables composable without starting over.", eyebrow="Section 03", logo_path=logo_white)
    slide_content(prs, "The VTEX composable architecture", bullets=[
        "Native headless via VTEX IO and FastStore — deploy your storefront independently",
        "Single OMS managing B2C, B2B, Marketplace, and Omnichannel from one backend",
        "500+ native APIs covering catalog, checkout, pricing, promotions, logistics",
        "VTEX App Store for certified third-party integrations — no custom glue code",
        "Managed SaaS infrastructure: automatic scaling, global CDN, zero downtime deploys",
    ], logo_path=logo_rebel)
    slide_content(prs, "Three migration paths — pick your pace", bullets=[
        "Lift-and-shift: migrate the full stack, keep existing integrations, optimize later",
        "Headless-first: decouple the frontend now, migrate backend capabilities gradually",
        "Capability-by-capability: replace only what's broken — checkout, catalog, OMS — independently",
    ], note="All three paths share the same VTEX backend — no data migration required between phases.",
       logo_path=logo_rebel)
    slide_quote(prs,
                quote="We went from deploying once a quarter to deploying twice a week. "
                      "The frontend team stopped waiting for backend releases.",
                attribution="CTO, Global Fashion Retailer — VTEX customer", logo_path=logo_rebel)
    slide_stat(prs, title="What merchants achieve after moving to composable on VTEX", stats=[
        ("2\u00d7",    "Faster frontend\ndeployments"),
        ("-35%",  "Reduction in\ntotal platform cost"),
        ("99.99%", "Uptime SLA\nacross Black Friday peaks"),
    ], logo_path=logo_white)
    slide_closing(prs, "Where do you want to be\nin 18 months?", cta="vtex.com/enterprise", logo_path=logo_white)

    prs.save(output_path)
    print(f"Deck saved: {output_path} ({len(prs.slides)} slides)")


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Generate VTEX Composable Commerce deck')
    parser.add_argument('--output',     default='vtex_composable_commerce.pptx')
    parser.add_argument('--logo-rebel', default=None, help='Rebel Pink logo PNG (light slides)')
    parser.add_argument('--logo-white', default=None, help='White logo PNG (dark/pink slides)')
    args = parser.parse_args()
    build_deck(args.output, args.logo_rebel, args.logo_white)
